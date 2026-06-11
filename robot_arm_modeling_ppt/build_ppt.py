# -*- coding: utf-8 -*-
"""Build the 3-page (plus cover) PPT: arm grasp dynamics -> state space -> transfer function."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")

NAVY   = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x3A, 0x7C, 0xA5)
DEEP   = RGBColor(0x2E, 0x5E, 0x8C)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xEA, 0xF1, 0xF6)
CODEBG = RGBColor(0xF4, 0xF6, 0xF8)
GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
ORANGE = RGBColor(0xB0, 0x5A, 0x1A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

EA_FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def style(run, size=14, bold=False, color=NAVY, latin="Calibri", ea=EA_FONT, italic=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = latin
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn("a:ea"))
    if ea_el is None:
        ea_el = etree.SubElement(rPr, qn("a:ea"))
    ea_el.set("typeface", ea)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, runs, align=PP_ALIGN.LEFT, space_before=0, space_after=6, line=None, first=False):
    """runs: list of (text, kwargs) tuples."""
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    for txt, kw in runs:
        r = p.add_run(); r.text = txt; style(r, **kw)
    return p


def rect(slide, x, y, w, h, fill, line_color=None, line_w=0.75, shadow=False, round_=False, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def pic(slide, path, x=None, y=None, h=None, w=None, center_in=None, max_w=None):
    """center_in=(cx, cw): horizontally center within that span."""
    im = Image.open(path); ar = im.width / im.height
    if h is not None:
        w = h * ar
    elif w is not None:
        h = w / ar
    if max_w is not None and w > max_w:
        w = max_w; h = w / ar
    if center_in is not None:
        x = center_in[0] + (center_in[1] - w) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return w, h


def header(slide, title, num):
    rect(slide, 0, 0, 13.333, 0.06, ACCENT)             # hairline top
    rect(slide, 0.55, 0.40, 0.13, 0.42, ACCENT)         # accent tab
    tf = textbox(slide, 0.84, 0.30, 11.6, 0.62)
    para(tf, [(title, dict(size=23, bold=True, color=NAVY))], first=True, space_after=0)
    rect(slide, 0.55, 1.02, 12.23, 0.012, RGBColor(0xD5, 0xDD, 0xE5))
    # footer
    tf = textbox(slide, 0.55, 7.10, 9.5, 0.3)
    para(tf, [("机械臂抓取任务建模：动力学方程 — 状态空间 — 传递函数",
               dict(size=9, color=GRAY))], first=True, space_after=0)
    tf = textbox(slide, 12.3, 7.10, 0.5, 0.3)
    para(tf, [(str(num), dict(size=10, color=GRAY))], align=PP_ALIGN.RIGHT, first=True, space_after=0)


def bullet(tf, segs, size=14, gap=6, sub=False, line=1.12, first=False):
    """segs: list of (text, extra_kwargs)"""
    mark = "–  " if sub else "▪  "
    runs = [(mark, dict(size=size, color=ACCENT if not sub else GRAY, bold=not sub))]
    for t, kw in segs:
        base = dict(size=size, color=NAVY); base.update(kw)
        runs.append((t, base))
    return para(tf, runs, space_after=gap, line=line, first=first)


# ===================== Slide 1: cover =====================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 0.10, ACCENT)
rect(s, 0, 7.40, 13.333, 0.10, ACCENT)

tf = textbox(s, 1.0, 1.55, 11.33, 1.0)
para(tf, [("机械臂抓取任务的动态系统建模", dict(size=36, bold=True, color=NAVY))],
     align=PP_ALIGN.CENTER, first=True, space_after=10)
tf = textbox(s, 1.0, 2.55, 11.33, 0.5)
para(tf, [("动力学方程  →  状态空间模型  →  传递函数    |    MATLAB 实现",
           dict(size=17, color=GRAY))], align=PP_ALIGN.CENTER, first=True)

# flow boxes
bx_w, bx_h, by = 3.30, 1.30, 3.85
xs = [0.95, 5.02, 9.09]
titles = ["① 动力学建模", "② 状态空间变换", "③ 传递函数"]
subs = ["拉格朗日方程（含抓取负载）", "平衡点线性化 → (A, B, C, D)", "G(s) = C(sI−A)⁻¹B + D"]
for x, t, sub in zip(xs, titles, subs):
    rect(s, x, by, bx_w, bx_h, LIGHT, line_color=ACCENT, line_w=1.2, round_=True, radius=0.14)
    tf = textbox(s, x + 0.15, by + 0.20, bx_w - 0.3, 0.5)
    para(tf, [(t, dict(size=17, bold=True, color=DEEP))], align=PP_ALIGN.CENTER, first=True, space_after=4)
    tf = textbox(s, x + 0.15, by + 0.72, bx_w - 0.3, 0.45)
    para(tf, [(sub, dict(size=11.5, color=GRAY))], align=PP_ALIGN.CENTER, first=True, space_after=0)
for x in (4.33, 8.40):
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(by + 0.47), Inches(0.62), Inches(0.36))
    ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT; ar.line.fill.background(); ar.shadow.inherit = False

tf = textbox(s, 1.0, 6.35, 11.33, 0.4)
para(tf, [("二连杆平面机械臂 · 末端抓取负载 · 2026 年 6 月", dict(size=13, color=GRAY))],
     align=PP_ALIGN.CENTER, first=True)

# ===================== Slide 2: dynamics modeling =====================
s = prs.slides.add_slide(BLANK)
header(s, "一、抓取任务的动力学建模（拉格朗日方法）", 1)

LX, LW = 0.55, 7.45
tf = textbox(s, LX, 1.22, LW, 1.30)
bullet(tf, [("对象：", dict(bold=True)),
            ("平面二连杆机械臂，关节力矩 τ1、τ2 驱动；末端抓取质量为 mp 的负载（视为与连杆 2 固连的点质量）。", {})],
       first=True, size=13.5)
bullet(tf, [("建模：", dict(bold=True)),
            ("取广义坐标 q = [q1, q2]ᵀ，写出系统动能 T 与势能 V（均含负载项），代入拉格朗日方程：", {})],
       size=13.5)
pic(s, f"{A}/eq_lagrange.png", y=2.42, h=0.66, center_in=(LX, LW))

tf = textbox(s, LX, 3.22, LW, 0.4)
bullet(tf, [("得到抓取系统的", {}), ("动力学方程", dict(bold=True, color=DEEP)),
            ("（非线性二阶动态系统模型）：", {})], first=True, size=13.5)
rect(s, LX + 0.45, 3.70, LW - 0.9, 0.78, LIGHT, line_color=ACCENT, line_w=1.0, round_=True, radius=0.18)
pic(s, f"{A}/eq_dynamics.png", y=3.86, h=0.46, center_in=(LX, LW))

tf = textbox(s, LX, 4.66, LW, 1.75)
bullet(tf, [("M(q)∈ℝ²ˣ²：惯性矩阵，对称正定，含负载 mp 的附加惯量项", {})], sub=True, size=12.5, gap=3, first=True)
bullet(tf, [("C 项：科氏力 / 离心力项；  G(q)：重力矩项；  Fv：关节粘性摩擦阻尼", {})], sub=True, size=12.5, gap=3)
bullet(tf, [("τ：关节驱动力矩（控制输入）；  q：关节角（被控量）", {})], sub=True, size=12.5, gap=8)
bullet(tf, [("抓取的影响：", dict(bold=True)),
            ("负载使 M(q) 与 G(q) 中出现含 mp 的项（等效惯量与重力矩增大），模型随抓取状态改变 —— 这正是需要建模分析的原因。", {})],
       size=13.5)

tf = textbox(s, LX, 6.45, LW, 0.55)
para(tf, [("本例参数：m1=2.0, m2=1.0, mp=0.5 kg；l1=0.5, l2=0.4 m；lc1=0.25, lc2=0.2 m；"
           "I1=0.05, I2=0.02 kg·m²；Fv=diag(0.5, 0.3) N·m·s/rad",
           dict(size=10, color=GRAY))], first=True, line=1.15)

RX, RW = 8.30, 4.45
rect(s, RX - 0.10, 1.20, RW + 0.20, 5.78, WHITE, line_color=RGBColor(0xD5, 0xDD, 0xE5), line_w=1.0, round_=True, radius=0.05)
pic(s, f"{A}/arm_diagram.png", y=1.38, h=4.95, center_in=(RX, RW))
tf = textbox(s, RX, 6.42, RW, 0.45)
para(tf, [("图：二连杆机械臂抓取负载示意（基座悬挂式）", dict(size=10.5, color=GRAY))],
     align=PP_ALIGN.CENTER, first=True)

# ===================== Slide 3: state space =====================
s = prs.slides.add_slide(BLANK)
header(s, "二、平衡点线性化与状态空间模型", 2)

LX, LW = 0.55, 7.55
tf = textbox(s, LX, 1.22, LW, 0.85)
bullet(tf, [("选取平衡点（工作点）：", dict(bold=True)),
            ("悬垂抓取位形 q0 = [−90°, 0°]ᵀ、速度为零；该位形下重力矩 G(q0) = 0，故平衡输入 τ0 = 0。", {})],
       first=True, size=13.5)
bullet(tf, [("小偏差线性化：", dict(bold=True)),
            ("令 δq = q − q0，速度乘积项（C 项）为二阶小量、舍去，得线性化方程：", {})], size=13.5)
pic(s, f"{A}/eq_linear.png", y=2.46, h=0.60, center_in=(LX, LW), max_w=LW - 0.15)

tf = textbox(s, LX, 3.22, LW, 0.4)
bullet(tf, [("选状态向量 x（角度 + 角速度）、输入 u、输出 y：", {})], first=True, size=13.5)
pic(s, f"{A}/eq_state.png", y=3.66, h=0.40, center_in=(LX, LW))

tf = textbox(s, LX, 4.28, LW, 0.4)
bullet(tf, [("改写为状态空间标准形，得到", {}), ("状态空间矩阵 (A, B, C, D)", dict(bold=True, color=DEEP)),
            ("：", {})], first=True, size=13.5)
pic(s, f"{A}/eq_ss.png", y=4.74, h=0.42, center_in=(LX, LW))
rect(s, LX + 0.10, 5.36, LW - 0.2, 1.32, LIGHT, line_color=ACCENT, line_w=1.0, round_=True, radius=0.10)
pic(s, f"{A}/eq_blocks.png", y=5.52, h=1.00, center_in=(LX, LW))

RX, RW = 8.42, 4.35
rect(s, RX - 0.12, 1.22, RW + 0.24, 5.46, CODEBG, line_color=RGBColor(0xD5, 0xDD, 0xE5), line_w=1.0, round_=True, radius=0.06)
tf = textbox(s, RX, 1.42, RW, 0.4)
para(tf, [("代入数值（含 0.5 kg 负载）", dict(size=13, bold=True, color=DEEP))], first=True)
mono = [
    ("M0 = [1.090  0.340", NAVY),
    ("      0.340  0.140]", NAVY),
    ("Kg = [16.187  3.924", NAVY),
    ("       3.924  3.924]", NAVY),
    ("", NAVY),
    ("A =", DEEP),
    ("     0      0      1.000   0", NAVY),
    ("     0      0      0       1.000", NAVY),
    ("  -25.188  21.211 -1.892   2.757", NAVY),
    ("   33.142 -79.541  4.595  -8.838", NAVY),
    ("", NAVY),
    ("B =", DEEP),
    ("    0       0", NAVY),
    ("    0       0", NAVY),
    ("    3.784  -9.189", NAVY),
    ("   -9.189  29.459", NAVY),
]
tf = textbox(s, RX, 1.86, RW, 4.2)
for i, (ln, col) in enumerate(mono):
    para(tf, [(ln if ln else " ", dict(size=11, color=col, latin="Consolas", bold=ln.endswith("=")))],
         first=(i == 0), space_after=1.5, line=1.0)
tf = textbox(s, RX, 6.16, RW, 0.5)
para(tf, [("C = [I₂  0₂]，D = 0；A 的特征值即系统极点（下页）",
           dict(size=10.5, color=GRAY))], first=True, line=1.15)

# ===================== Slide 4: transfer function + MATLAB =====================
s = prs.slides.add_slide(BLANK)
header(s, "三、状态空间 → 传递函数：MATLAB 实现与结果", 3)

# left: code
CX, CW = 0.55, 6.30
rect(s, CX - 0.06, 1.18, CW + 0.12, 5.72, CODEBG, line_color=RGBColor(0xD5, 0xDD, 0xE5), line_w=1.0, round_=True, radius=0.05)
code = [
    "%% 机械臂抓取: 动力学 -> 状态空间 -> 传递函数",
    "M  = [1.090 0.340; 0.340 0.140];    % 惯性矩阵 M(q0), 含负载",
    "Kg = [16.187 3.924; 3.924 3.924];   % 重力刚度 dG/dq|q0",
    "Fv = diag([0.5 0.3]);               % 关节粘性阻尼",
    "n  = 2;",
    "",
    "%% 状态空间模型: x = [dq; dq'], u = dtau, y = dq",
    "A = [zeros(n), eye(n); -M\\Kg, -M\\Fv];",
    "B = [zeros(n); inv(M)];",
    "C = [eye(n), zeros(n)];   D = zeros(n);",
    "sys = ss(A, B, C, D);",
    "",
    "%% 传递函数矩阵 G(s) = C(sI-A)^(-1)B + D",
    "G   = tf(sys);",
    "G11 = minreal(G(1,1))               % 通道: dtau1 -> dq1",
    "[num, den] = ss2tf(A, B, C, D, 1);  % 等价做法(第1输入)",
    "",
    "%% 验证与分析",
    "p = pole(sys)                       % 极点 -> 稳定性",
    "dcgain(G11)                         % 直流增益 = inv(Kg)(1,1)",
    "step(G11), grid on                  % 阶跃响应",
    "bode(G11), grid on                  % 频率特性",
]
tf = textbox(s, CX + 0.14, 1.36, CW - 0.24, 5.4)
for i, ln in enumerate(code):
    if "%" in ln:
        cut = ln.index("%")
        if ln.strip().startswith("%%"):
            runs = [(ln, dict(size=10.5, color=GREEN, latin="Consolas", bold=True))]
        else:
            runs = [(ln[:cut], dict(size=10.5, color=NAVY, latin="Consolas")),
                    (ln[cut:], dict(size=10.5, color=GREEN, latin="Consolas"))]
    else:
        runs = [(ln if ln else " ", dict(size=10.5, color=NAVY, latin="Consolas"))]
    para(tf, runs, first=(i == 0), space_after=2.2, line=1.0)

# right column
RX, RW = 7.15, 5.65
tf = textbox(s, RX, 1.20, RW, 0.4)
bullet(tf, [("由状态空间矩阵求", {}), ("传递函数", dict(bold=True, color=DEEP)), ("（2×2 矩阵，MIMO）：", {})],
       first=True, size=13.5)
pic(s, f"{A}/eq_tf.png", y=1.70, h=0.50, center_in=(RX, RW), max_w=RW - 0.15)

tf = textbox(s, RX, 2.42, RW, 0.4)
bullet(tf, [("MATLAB 运行结果（取 δτ1 → δq1 通道）：", {})], first=True, size=13.5)
rect(s, RX + 0.05, 2.86, RW - 0.1, 1.02, LIGHT, line_color=ACCENT, line_w=1.0, round_=True, radius=0.12)
pic(s, f"{A}/eq_g11.png", y=3.06, h=0.68, center_in=(RX + 0.05, RW - 0.1), max_w=RW - 0.55)

tf = textbox(s, RX, 4.10, RW, 2.6)
bullet(tf, [("极点 pole(sys)（rad/s）：", dict(bold=True))], first=True, size=13, gap=2)
para(tf, [("    p = -0.196 ± 3.789j ,  -5.169 ± 7.976j",
           dict(size=12.5, color=DEEP, latin="Consolas", bold=True))], space_after=7)
bullet(tf, [("四个极点均位于左半平面 ⇒ 悬垂抓取平衡位形渐近稳定；主导极点 ωn ≈ 3.79 rad/s、ζ ≈ 0.05，呈欠阻尼摆动（阻尼弱）。", {})],
       sub=True, size=12.5, gap=8)
bullet(tf, [("直流增益 dcgain(G11) = 0.0815 rad/(N·m)，与 Kg 逆矩阵的 (1,1) 元素一致，模型自洽。", {})],
       sub=True, size=12.5, gap=8)
bullet(tf, [("用途：", dict(bold=True)),
            ("此 (A,B,C,D) 与 G(s) 即为控制设计的出发点 —— 可进一步设计 PID / LQR / 状态观测器，并分析负载 mp 变化对抓取动态特性的影响。", {})],
       size=13)
tf = textbox(s, RX, 6.62, RW, 0.4)
para(tf, [("完整可运行脚本：grasp_dynamics_ss_tf.m（由物理参数自动计算 M0、Kg 并绘制 step/bode）",
           dict(size=10, color=GRAY))], first=True)

prs.core_properties.title = "机械臂抓取任务的动态系统建模：动力学方程-状态空间-传递函数"
prs.core_properties.author = ""
out = os.path.join(HERE, "机械臂抓取任务建模_动力学_状态空间_传递函数.pptx")
prs.save(out)
print("saved", out)
